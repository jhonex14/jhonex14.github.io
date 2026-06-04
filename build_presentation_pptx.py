import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ----------------------------------------------------
    # COLOR PALETTE SETUP
    # ----------------------------------------------------
    # ConsulTime Theme Colors
    C_NAVY_DARK  = RGBColor(15, 23, 42)      # #0F172A
    C_NAVY_LIGHT = RGBColor(30, 58, 138)     # #1E3A8A
    C_BLUE       = RGBColor(59, 130, 246)     # #3B82F6
    C_EMERALD    = RGBColor(16, 185, 129)     # #10B981
    C_MINT       = RGBColor(52, 211, 153)     # #34D399
    C_BG_LIGHT   = RGBColor(248, 250, 252)    # #F8FAFC
    C_WHITE      = RGBColor(255, 255, 255)    # #FFFFFF
    C_TEXT_DARK  = RGBColor(31, 41, 55)      # #1F2937
    C_TEXT_MUTED = RGBColor(107, 114, 128)    # #6B7280
    C_GRAY_LIGHT = RGBColor(229, 231, 235)    # #E5E7EB
    C_RED_SOFT   = RGBColor(254, 242, 242)    # #FEF2F2
    C_RED        = RGBColor(239, 68, 68)      # #EF4444

    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # Helper function to add slide background
    def set_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper to draw colored rectangle shapes (cards, panels, accents)
    def draw_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=1):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(border_width)
        else:
            shape.line.fill.background() # No border
        return shape

    # Helper to add modern slide titles with primary & accent text styling
    def add_slide_header(slide, title_text, subtitle_text=None, is_dark_bg=False):
        # Draw top accent bar
        draw_rect(slide, Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.08), C_EMERALD if is_dark_bg else C_NAVY_LIGHT)
        
        # Add Header Textbox
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = C_WHITE if is_dark_bg else C_NAVY_LIGHT
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = "Arial"
            p2.font.size = Pt(14)
            p2.font.color.rgb = C_MINT if is_dark_bg else C_TEXT_MUTED

    # Helper to create styled bullet lists inside textboxes
    def create_textbox_list(slide, left, top, width, height, bullets, font_size=13, font_color=C_TEXT_DARK, bold_title_prefix=True):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        for idx, bullet in enumerate(bullets):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.font.name = "Arial"
            p.font.size = Pt(font_size)
            p.space_after = Pt(8)
            p.level = 0
            
            # Highlight bold prefix (e.g. "Main Task: explanation")
            if bold_title_prefix and ":" in bullet:
                parts = bullet.split(":", 1)
                run1 = p.add_run()
                run1.text = "• " + parts[0] + ":"
                run1.font.bold = True
                run1.font.color.rgb = font_color
                
                run2 = p.add_run()
                run2.text = parts[1]
                run2.font.bold = False
                run2.font.color.rgb = font_color
            else:
                p.text = "• " + bullet
                p.font.color.rgb = font_color
        return tb

    # Helper to draw a CSS-style browser mockup window to host UI screenshots/sketches
    def draw_browser_mockup(slide, left, top, width, height, title="ConsulTime Portal - Mockup"):
        # Browser Header bar
        header_h = Inches(0.4)
        draw_rect(slide, left, top, width, header_h, C_NAVY_LIGHT)
        # Browser Content Canvas
        draw_rect(slide, left, top + header_h, width, height - header_h, C_WHITE, C_GRAY_LIGHT, 1.5)
        
        # Draw Window Control Dots (Minimize, Maximize, Close)
        for i, color in enumerate([RGBColor(239, 68, 68), RGBColor(245, 158, 11), RGBColor(16, 185, 129)]):
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.12 + 0.18*i), top + Inches(0.12), Inches(0.14), Inches(0.14))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()
            
        # Draw Title
        tb = slide.shapes.add_textbox(left + Inches(0.8), top + Inches(0.04), width - Inches(1.5), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = C_WHITE

    # ----------------------------------------------------
    # SLIDE 1: TITLE SLIDE (Dark Navy)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_NAVY_DARK)
    
    # Background glowing accents
    draw_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), C_EMERALD)
    draw_rect(slide, Inches(13.18), Inches(0), Inches(0.15), Inches(7.5), C_BLUE)
    
    # Big title box
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "CONSULTIME"
    p1.font.name = "Arial"
    p1.font.size = Pt(56)
    p1.font.bold = True
    p1.font.color.rgb = C_MINT
    p1.alignment = PP_ALIGN.LEFT
    p1.space_after = Pt(14)
    
    p2 = tf.add_paragraph()
    p2.text = "A Web and Mobile-Based Academic Consultation and Appointment Scheduling System with Real-Time Authentication"
    p2.font.name = "Arial"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE
    p2.alignment = PP_ALIGN.LEFT
    p2.space_after = Pt(45)
    
    p3 = tf.add_paragraph()
    p3.text = "Researchers: Khyle Bulawan, Carlvyn Bajala, Jake Pinggol, Henan Oliveros"
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = C_WHITE
    
    p4 = tf.add_paragraph()
    p4.text = "Institution: Philippine College of Technology  |  Date: June 2026"
    p4.font.name = "Arial"
    p4.font.size = Pt(13)
    p4.font.color.rgb = C_BLUE

    # ----------------------------------------------------
    # SLIDE 2: PROBLEM STATEMENT (Light Gray - Split Layout)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "PROBLEM STATEMENT", "Challenges in manual consultation scheduling and communication at PCT")
    
    # Left Column: The Baseline Bottleneck
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(4.0), Inches(5.0), C_NAVY_DARK)
    tb_left = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(3.4), Inches(4.2))
    tf_left = tb_left.text_frame
    tf_left.word_wrap = True
    p_left1 = tf_left.paragraphs[0]
    p_left1.text = "The Scheduling Bottleneck"
    p_left1.font.bold = True
    p_left1.font.size = Pt(18)
    p_left1.font.color.rgb = C_MINT
    p_left1.space_after = Pt(14)
    p_left2 = tf_left.add_paragraph()
    p_left2.text = "Traditional consultation logistics rely on physical bullet-boards, uncoordinated email channels, and manual logs.\n\nThis setup generates coordination friction, wastes academic hours, and leads to high student-faculty mismatch rates."
    p_left2.font.size = Pt(13.5)
    p_left2.font.color.rgb = C_WHITE
    p_left2.space_before = Pt(10)
    
    # Right Column: Four Key Pain Points (Grid layout)
    problems = [
        "Uncoordinated Slots: Faculty availability lists are offline and static, causing double-bookings.",
        "Communication Lag: Users lack instant notification systems for sudden schedule modifications.",
        "Authentication Risks: Registrations lack administrative verification pipelines, allowing fake accounts.",
        "Poor Mobile Adoption: Portals lack mobile wrappers or touch-friendly apps, reducing student accessibility."
    ]
    for idx, prob in enumerate(problems):
        row = idx // 2
        col = idx % 2
        x = Inches(4.85 + col * 4.0)
        y = Inches(1.8 + row * 2.5)
        
        # Soft Red accent line for problem cards
        draw_rect(slide, x, y, Inches(3.7), Inches(2.3), C_WHITE, C_GRAY_LIGHT)
        draw_rect(slide, x, y, Inches(0.08), Inches(2.3), C_RED)
        
        create_textbox_list(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.3), Inches(1.9), [prob], font_size=12.5)

    # ----------------------------------------------------
    # SLIDE 3: OBJECTIVES OF THE STUDY (Light Gray)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "OBJECTIVES", "Optimizing academic collaboration through digital transformation")
    
    # Left Box: Primary Objective
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(4.5), Inches(5.0), C_NAVY_LIGHT)
    tb_obj = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(3.9), Inches(4.4))
    tf_obj = tb_obj.text_frame
    tf_obj.word_wrap = True
    p_obj1 = tf_obj.paragraphs[0]
    p_obj1.text = "Primary Objective"
    p_obj1.font.bold = True
    p_obj1.font.size = Pt(20)
    p_obj1.font.color.rgb = C_MINT
    p_obj1.space_after = Pt(20)
    p_obj2 = tf_obj.add_paragraph()
    p_obj2.text = "To design, develop, and evaluate ConsulTime—a real-time, authenticated web and mobile scheduling system that bridges the coordination gap between students and faculty members at PCT."
    p_obj2.font.size = Pt(14)
    p_obj2.font.color.rgb = C_WHITE
    p_obj2.space_before = Pt(10)
    
    # Right Box: Specific Objectives
    draw_rect(slide, Inches(5.3), Inches(1.8), Inches(7.5), Inches(5.0), C_WHITE, C_GRAY_LIGHT)
    draw_rect(slide, Inches(5.3), Inches(1.8), Inches(0.12), Inches(5.0), C_EMERALD)
    
    tb_spec = slide.shapes.add_textbox(Inches(5.7), Inches(2.0), Inches(6.8), Inches(4.6))
    tf_spec = tb_spec.text_frame
    tf_spec.word_wrap = True
    p_spec_title = tf_spec.paragraphs[0]
    p_spec_title.text = "Specific Project Targets:"
    p_spec_title.font.bold = True
    p_spec_title.font.size = Pt(17)
    p_spec_title.font.color.rgb = C_NAVY_LIGHT
    p_spec_title.space_after = Pt(12)
    
    spec_bullets = [
        "Dashboard Interfaces: Establish dedicated views for Students, Faculty, and Admins.",
        "Real-Time Manager: Allow faculty to open, update, and close scheduling slots dynamically.",
        "Unified Booking: Enable students to search by course or department and book open slots instantly.",
        "Capacitor Wrapper: Compile the web portal into a lightweight native Android APK.",
        "Vetting Pipeline: Secure the platform by letting Admins approve faculty credentials.",
        "Live Version Telemetry: Sync and prompt clients to update local assets for new app versions."
    ]
    
    create_textbox_list(slide, Inches(5.7), Inches(2.5), Inches(6.8), Inches(4.1), spec_bullets, font_size=12.5)

    # ----------------------------------------------------
    # SLIDE 4: SYSTEM ARCHITECTURE (Dark Navy)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_NAVY_DARK)
    add_slide_header(slide, "SYSTEM ARCHITECTURE", "High-level integration and data flow layout of the system", is_dark_bg=True)
    
    # Drawing flowchart steps:
    # Step 1: User Interfaces
    draw_rect(slide, Inches(0.8), Inches(2.5), Inches(3.2), Inches(3.5), C_NAVY_LIGHT, C_BLUE, 2)
    tb_arch1 = slide.shapes.add_textbox(Inches(1.0), Inches(2.7), Inches(2.8), Inches(3.1))
    tf_arch1 = tb_arch1.text_frame
    tf_arch1.word_wrap = True
    tf_arch1.paragraphs[0].text = "CLIENT TIER"
    tf_arch1.paragraphs[0].font.bold = True
    tf_arch1.paragraphs[0].font.size = Pt(15)
    tf_arch1.paragraphs[0].font.color.rgb = C_MINT
    tf_arch1.paragraphs[0].space_after = Pt(10)
    create_textbox_list(slide, Inches(0.9), Inches(3.2), Inches(3.0), Inches(2.7), [
        "Web Application (HTML5/Vanilla CSS/ES6 JS)",
        "Mobile App: Capacitor Native Android Wrapper (.APK)",
        "Local Cache: Service Workers (sw.js) for fast assets reload"
    ], font_size=11.5, font_color=C_WHITE, bold_title_prefix=False)

    # Arrow 1 to 2
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.2), Inches(4.0), Inches(0.7), Inches(0.5))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = C_BLUE
    arrow1.line.fill.background()

    # Step 2: Cloud Backend
    draw_rect(slide, Inches(5.1), Inches(2.5), Inches(3.2), Inches(3.5), C_NAVY_LIGHT, C_EMERALD, 2)
    tb_arch2 = slide.shapes.add_textbox(Inches(5.3), Inches(2.7), Inches(2.8), Inches(3.1))
    tf_arch2 = tb_arch2.text_frame
    tf_arch2.word_wrap = True
    tf_arch2.paragraphs[0].text = "DATABASE & AUTH TIER"
    tf_arch2.paragraphs[0].font.bold = True
    tf_arch2.paragraphs[0].font.size = Pt(15)
    tf_arch2.paragraphs[0].font.color.rgb = C_MINT
    tf_arch2.paragraphs[0].space_after = Pt(10)
    create_textbox_list(slide, Inches(5.2), Inches(3.2), Inches(3.0), Inches(2.7), [
        "Supabase BaaS PostgreSQL cloud storage instance",
        "Row Level Security (RLS) policies for user data security",
        "Real-Time Listeners for database slot synchronization"
    ], font_size=11.5, font_color=C_WHITE, bold_title_prefix=False)

    # Arrow 2 to 3
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.5), Inches(4.0), Inches(0.7), Inches(0.5))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = C_EMERALD
    arrow2.line.fill.background()

    # Step 3: Admin Vetting Pipeline
    draw_rect(slide, Inches(9.4), Inches(2.5), Inches(3.2), Inches(3.5), C_NAVY_LIGHT, C_BLUE, 2)
    tb_arch3 = slide.shapes.add_textbox(Inches(9.6), Inches(2.7), Inches(2.8), Inches(3.1))
    tf_arch3 = tb_arch3.text_frame
    tf_arch3.word_wrap = True
    tf_arch3.paragraphs[0].text = "ADMINISTRATIVE TIER"
    tf_arch3.paragraphs[0].font.bold = True
    tf_arch3.paragraphs[0].font.size = Pt(15)
    tf_arch3.paragraphs[0].font.color.rgb = C_MINT
    tf_arch3.paragraphs[0].space_after = Pt(10)
    create_textbox_list(slide, Inches(9.5), Inches(3.2), Inches(3.0), Inches(2.7), [
        "Vetting Panel: Checks and approves faculty registrations",
        "Deployment: Publishes Direct Download Android links",
        "Version Control: Updates version.json to sync app clients"
    ], font_size=11.5, font_color=C_WHITE, bold_title_prefix=False)

    # ----------------------------------------------------
    # SLIDE 5: TECHNOLOGIES USED (Light Gray)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "TECHNOLOGIES USED", "Modern and robust tech stack powering ConsulTime")
    
    tech_stack = [
        ("HTML5 & Vanilla CSS3", "Ensures high styling flexibility, responsive mobile layout, and zero screen-flicker during page render."),
        ("ES6 JavaScript", "Core programming logic managing dynamic dashboard updates and browser event hooks."),
        ("Supabase Cloud Database", "PostgreSQL database tier managing secure users profiles, open slots, and slot bookings."),
        ("Supabase Authentication", "Implements session caching, user profiles mapping, and custom metadata."),
        ("Capacitor Mobile", "Packages the web portal into a native, lightweight Android application wrapper (.APK)."),
        ("Service Worker API", "Handles asset caching (manifest.json, sw.js) to guarantee offline availability.")
    ]
    
    for idx, tech in enumerate(tech_stack):
        row = idx // 3
        col = idx % 3
        x = Inches(0.5 + col * 4.2)
        y = Inches(1.8 + row * 2.6)
        
        draw_rect(slide, x, y, Inches(3.9), Inches(2.3), C_WHITE, C_GRAY_LIGHT)
        draw_rect(slide, x, y, Inches(3.9), Inches(0.08), C_BLUE if row == 0 else C_EMERALD)
        
        tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(3.5), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = tech[0]
        p.font.bold = True
        p.font.size = Pt(15)
        p.font.color.rgb = C_NAVY_LIGHT
        p.space_after = Pt(10)
        
        p2 = tf.add_paragraph()
        p2.text = tech[1]
        p2.font.size = Pt(12)
        p2.font.color.rgb = C_TEXT_DARK

    # ----------------------------------------------------
    # SLIDE 6-11: UI SHOWCASE MODULES (Mockups and Text descriptions)
    # ----------------------------------------------------
    # Helper to generate standard mockup showcase slides
    def add_mockup_showcase_slide(title_header, desc_title, desc_bullets, mockup_label, mockup_details):
        s = prs.slides.add_slide(blank_layout)
        set_background(s, C_BG_LIGHT)
        add_slide_header(s, title_header, "ConsulTime platform frontend showcase")
        
        # Left Side: Description card
        draw_rect(s, Inches(0.5), Inches(1.8), Inches(4.5), Inches(5.0), C_WHITE, C_GRAY_LIGHT)
        draw_rect(s, Inches(0.5), Inches(1.8), Inches(4.5), Inches(0.12), C_NAVY_LIGHT)
        
        tb_desc = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(3.9), Inches(4.3))
        tf_desc = tb_desc.text_frame
        tf_desc.word_wrap = True
        tf_desc.paragraphs[0].text = desc_title
        tf_desc.paragraphs[0].font.bold = True
        tf_desc.paragraphs[0].font.size = Pt(18)
        tf_desc.paragraphs[0].font.color.rgb = C_NAVY_LIGHT
        tf_desc.paragraphs[0].space_after = Pt(14)
        
        create_textbox_list(s, Inches(0.8), Inches(2.7), Inches(3.9), Inches(3.8), desc_bullets, font_size=12.5)
        
        # Right Side: Browser Mockup hosting a custom schema representation
        mx, my, mw, mh = Inches(5.4), Inches(1.8), Inches(7.4), Inches(5.0)
        draw_browser_mockup(s, mx, my, mw, mh, title=f"ConsulTime - {mockup_label}")
        
        # Mockup elements layout
        cx = mx + Inches(0.5)
        cy = my + Inches(0.8)
        cw = mw - Inches(1.0)
        ch = mh - Inches(1.2)
        
        # Draw background inside canvas
        draw_rect(s, cx, cy, cw, ch, C_BG_LIGHT, C_GRAY_LIGHT)
        
        # Add visual components representing mockup page layout
        y_offset = cy + Inches(0.2)
        for row_label, badge, badge_color in mockup_details:
            # Row container card
            draw_rect(s, cx + Inches(0.2), y_offset, cw - Inches(0.4), Inches(0.7), C_WHITE, C_GRAY_LIGHT)
            
            # Row icon placeholder
            draw_rect(s, cx + Inches(0.35), y_offset + Inches(0.15), Inches(0.4), Inches(0.4), C_BLUE)
            
            # Row labels
            tb_lbl = s.shapes.add_textbox(cx + Inches(0.9), y_offset + Inches(0.1), cw - Inches(2.5), Inches(0.5))
            tf_lbl = tb_lbl.text_frame
            tf_lbl.word_wrap = True
            p_lbl = tf_lbl.paragraphs[0]
            p_lbl.text = row_label
            p_lbl.font.bold = True
            p_lbl.font.size = Pt(11)
            p_lbl.font.color.rgb = C_NAVY_DARK
            
            # Badge box
            draw_rect(s, cx + cw - Inches(1.8), y_offset + Inches(0.18), Inches(1.4), Inches(0.34), badge_color)
            tb_bdg = s.shapes.add_textbox(cx + cw - Inches(1.8), y_offset + Inches(0.15), Inches(1.4), Inches(0.34))
            p_bdg = tb_bdg.text_frame.paragraphs[0]
            p_bdg.text = badge
            p_bdg.alignment = PP_ALIGN.CENTER
            p_bdg.font.name = "Arial"
            p_bdg.font.size = Pt(9.5)
            p_bdg.font.bold = True
            p_bdg.font.color.rgb = C_WHITE
            
            y_offset += Inches(0.8)

    # Slide 6: Login Module
    add_mockup_showcase_slide(
        "LOGIN MODULE",
        "Secure User Credentials",
        [
            "Session Management: Leverages Supabase cache authentication to bypass loading screens.",
            "Role Routing: Auto-detects user role and redirects them to Student, Faculty, or Admin dashboard.",
            "Zero Screen Flicker: Injects verify check scripts before loading dashboard UI layouts.",
            "SSL Protocols: Encrypts login data streams via secure cloud sockets."
        ],
        "Authentication Login Screen",
        [
            ("Login Form Layout: input credentials", "Standard Form", C_NAVY_LIGHT),
            ("Role Selector: Select Student / Faculty", "User Selector", C_BLUE),
            ("System Redirect Console: Verify session", "Auto Redirect", C_EMERALD)
        ]
    )

    # Slide 7: Registration Module
    add_mockup_showcase_slide(
        "REGISTRATION MODULE",
        "University Sign-Up Pipeline",
        [
            "Sign-Up Validations: Filters email extensions, allowing only authorized student domain credentials.",
            "Credentials Archive: Prompts faculty users to upload valid institutional verification credentials.",
            "Inactive State: Freezes faculty accounts in a secure database table until admin vetting is complete.",
            "Row Security Sync: Protects private registration database files from public access."
        ],
        "Account Sign-Up & Register Console",
        [
            ("Register Form Layout: Input name, email, credentials", "Form Validator", C_NAVY_LIGHT),
            ("Upload Faculty Credentials: PDF/Image Upload", "Upload Vault", C_BLUE),
            ("Submit Pipeline: Pending Admin Activation", "Vetting Buffer", C_RED)
        ]
    )

    # Slide 8: Student Dashboard
    add_mockup_showcase_slide(
        "STUDENT DASHBOARD",
        "Appointment Control Console",
        [
            "Active Appointments: Lists ongoing, confirmed, and previous consultation sessions.",
            "Notifications Panel: Highlights instant alerts regarding scheduling shifts.",
            "Direct APK Download: Features prominent layout banners for mobile package downloads.",
            "UI Responsive Structure: Scales layout grids to fit standard smartphone viewports."
        ],
        "Student Scheduling Workspace",
        [
            ("Consultation List View: Syncing upcoming tasks", "Upcoming Slots", C_NAVY_LIGHT),
            ("Direct Mobile Access: Download Android APK", "Direct Link", C_BLUE),
            ("Version Warning Panel: Alert user of new changes", "System Check", C_EMERALD)
        ]
    )

    # Slide 9: Book Appointment
    add_mockup_showcase_slide(
        "BOOK APPOINTMENT MODULE",
        "Seamless Consultation Pipeline",
        [
            "Search Filters: Allows students to query faculty lists by name, course, or department.",
            "Dynamic Calendars: Synchronizes available slots from real-time Supabase calendars.",
            "Interactive Picking: Click to highlight preferred hour blocks and fill out meeting details.",
            "Row-Level Lock: Prevents other students from selecting already booked appointments."
        ],
        "Schedule Booking Wizard",
        [
            ("Search Input Panel: Filter faculty listings", "Live Filter", C_NAVY_LIGHT),
            ("Faculty Calendar Interface: Highlight available hours", "Time Slots", C_BLUE),
            ("Submit Reservation: RLS database locking", "Secure Booking", C_EMERALD)
        ]
    )

    # Slide 10: Faculty Dashboard
    add_mockup_showcase_slide(
        "FACULTY DASHBOARD",
        "Availability Blocks Manager",
        [
            "Calendar Controls: Allows faculty to define custom hours and block off busy dates.",
            "Status Workflow: Review booking requests and toggle status (Approve, Cancel, Complete).",
            "Automatic Conflicts Lock: Prevents booking overlapping slots, protecting research time.",
            "Activity Logs: Displays summary metrics of completed student advisory interactions."
        ],
        "Faculty Workspace Console",
        [
            ("Availability Slots Scheduler: Block off hours", "Slots Builder", C_NAVY_LIGHT),
            ("Pending Appointments Pipeline: Review requests", "Vetting Panel", C_BLUE),
            ("Update Status Tool: Approved / Cancelled / Done", "Status Badge", C_EMERALD)
        ]
    )

    # Slide 11: Admin Dashboard
    add_mockup_showcase_slide(
        "ADMIN DASHBOARD",
        "Platform Governance Panel",
        [
            "Vetting Console: Lists all pending faculty profiles with institutional uploads.",
            "Approval Action: Toggle profile activation, which instantly syncs user table permissions.",
            "System Logs Monitor: Tracks total active schedules, user roles, and database latency.",
            "Version Checker: Updates version.json parameters on the server to push client updates."
        ],
        "Administrator Control Console",
        [
            ("Faculty Account Vetting: Check verification files", "Verify Vault", C_NAVY_LIGHT),
            ("Telemetries Logger: System uptime, database latency", "Live Telemetry", C_BLUE),
            ("Update Release Tool: Version metadata manager", "version.json", C_EMERALD)
        ]
    )

    # ----------------------------------------------------
    # SLIDE 12: SYSTEM EVALUATION (Light Gray - Tables & Callout)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "SYSTEM EVALUATION", "Usability testing and query latency comparison results")
    
    # Left Column: Large Usability Metrics Callout Cards
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(4.5), Inches(2.3), C_WHITE, C_GRAY_LIGHT)
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(4.5), Inches(0.08), C_EMERALD)
    
    tb_usab = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(4.1), Inches(1.9))
    tf_usab = tb_usab.text_frame
    tf_usab.word_wrap = True
    p_usab = tf_usab.paragraphs[0]
    p_usab.text = "System Usability Score (SUS)"
    p_usab.font.bold = True
    p_usab.font.size = Pt(13)
    p_usab.font.color.rgb = C_TEXT_MUTED
    p_usab.space_after = Pt(4)
    p_usab2 = tf_usab.add_paragraph()
    p_usab2.text = "88.5 / 100"
    p_usab2.font.bold = True
    p_usab2.font.size = Pt(36)
    p_usab2.font.color.rgb = C_EMERALD
    p_usab3 = tf_usab.add_paragraph()
    p_usab3.text = "Average score evaluated by 30 students and 10 faculty members. Indicates excellent usability."
    p_usab3.font.size = Pt(11)
    p_usab3.font.color.rgb = C_TEXT_DARK
    
    draw_rect(slide, Inches(0.5), Inches(4.4), Inches(4.5), Inches(2.4), C_WHITE, C_GRAY_LIGHT)
    draw_rect(slide, Inches(0.5), Inches(4.4), Inches(4.5), Inches(0.08), C_BLUE)
    
    tb_perf = slide.shapes.add_textbox(Inches(0.7), Inches(4.6), Inches(4.1), Inches(2.0))
    tf_perf = tb_perf.text_frame
    tf_perf.word_wrap = True
    p_perf = tf_perf.paragraphs[0]
    p_perf.text = "Query Speed Optimization"
    p_perf.font.bold = True
    p_perf.font.size = Pt(13)
    p_perf.font.color.rgb = C_TEXT_MUTED
    p_perf.space_after = Pt(4)
    p_perf2 = tf_perf.add_paragraph()
    p_perf2.text = "99.7% Latency Reduction"
    p_perf2.font.bold = True
    p_perf2.font.size = Pt(28)
    p_perf2.font.color.rgb = C_BLUE
    p_perf3 = tf_perf.add_paragraph()
    p_perf3.text = "Average database query response times drop from 27.6 minutes (manual filing shelf search) to 3.6 seconds."
    p_perf3.font.size = Pt(11)
    p_perf3.font.color.rgb = C_TEXT_DARK

    # Right Column: Trial Data Table (Drawn inside PowerPoint)
    tb_tbl_lbl = slide.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(7.3), Inches(0.4))
    tb_tbl_lbl.text_frame.paragraphs[0].text = "Performance Telemetry: Manual vs. ConsulTime Search Latency"
    tb_tbl_lbl.text_frame.paragraphs[0].font.bold = True
    tb_tbl_lbl.text_frame.paragraphs[0].font.size = Pt(13.5)
    tb_tbl_lbl.text_frame.paragraphs[0].font.color.rgb = C_NAVY_LIGHT
    
    # Table layout
    rows = 7
    cols = 4
    left_x = Inches(5.5)
    top_y = Inches(2.3)
    width_w = Inches(7.3)
    height_h = Inches(4.5)
    
    table_shape = slide.shapes.add_table(rows, cols, left_x, top_y, width_w, height_h)
    tbl = table_shape.table
    
    # Set column widths
    tbl.columns[0].width = Inches(1.3)
    tbl.columns[1].width = Inches(2.2)
    tbl.columns[2].width = Inches(2.0)
    tbl.columns[3].width = Inches(1.8)
    
    # Headers
    headers = ["Trial", "Manual Search Time", "ConsulTime Search", "Time Saved"]
    for col_idx, header in enumerate(headers):
        cell = tbl.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_NAVY_LIGHT
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        
    # Table Data
    data = [
        ("Trial 1", "18.5 mins", "4.8 secs", "18.42 mins"),
        ("Trial 2", "34.0 mins", "3.0 secs", "33.95 mins"),
        ("Trial 3", "12.2 mins", "3.6 secs", "12.14 mins"),
        ("Trial 4", "45.0 mins", "4.2 secs", "44.93 mins"),
        ("Trial 5", "28.3 mins", "2.4 secs", "28.26 mins"),
        ("Average", "27.6 mins", "3.6 secs", "27.54 mins")
    ]
    
    for row_idx, row_data in enumerate(data):
        for col_idx, val in enumerate(row_data):
            cell = tbl.cell(row_idx + 1, col_idx)
            cell.text = val
            cell.fill.solid()
            # Highlight average row
            if row_idx == 5:
                cell.fill.fore_color.rgb = RGBColor(237, 245, 255) # Light blue accent
            else:
                cell.fill.fore_color.rgb = C_WHITE
            
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Arial"
            p.font.size = Pt(10.5)
            if row_idx == 5:
                p.font.bold = True
                p.font.color.rgb = C_NAVY_LIGHT
            else:
                p.font.color.rgb = C_TEXT_DARK

    # ----------------------------------------------------
    # SLIDE 13: DEVICE COMPATIBILITY (Light Gray)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "DEVICE COMPATIBILITY", "Ensuring accessible deployments on mobile and web viewports")
    
    compat_data = [
        ("Capacitor Mobile wrapper packaging (.APK)", [
            "Android compilation wrapper: Native package built to ensure quick execution speed on mobile systems.",
            "Touch-friendly panels: Viewports and tables scale correctly for single-finger operations.",
            "Light footprint: Minimizes device storage and memory usage compared to standard web browsers."
        ], C_BLUE),
        ("Online progressive Web application (PWA)", [
            "Responsive designs: Layout templates auto-fit standard notebooks, tables, and phone dimensions.",
            "Cached credentials: Uses Service Workers to save login configurations and prevent page refresh lags.",
            "Version telemetry checks: Auto-pings server version.json file to prompt quick cache invalidation."
        ], C_EMERALD)
    ]
    
    for idx, item in enumerate(compat_data):
        x = Inches(0.5 + idx * 6.3)
        y = Inches(1.8)
        
        draw_rect(slide, x, y, Inches(6.0), Inches(5.0), C_WHITE, C_GRAY_LIGHT)
        draw_rect(slide, x, y, Inches(6.0), Inches(0.12), item[2])
        
        tb = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.3), Inches(5.4), Inches(4.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = item[0]
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = C_NAVY_LIGHT
        p.space_after = Pt(18)
        
        create_textbox_list(slide, x + Inches(0.3), y + Inches(1.0), Inches(5.4), Inches(3.7), item[1], font_size=12.5)

    # ----------------------------------------------------
    # SLIDE 14: KEY CONTRIBUTIONS & INNOVATIONS (Dark Navy)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_NAVY_DARK)
    add_slide_header(slide, "KEY CONTRIBUTIONS", "Core system innovations and security benefits", is_dark_bg=True)
    
    contribs = [
        ("Vetting security pipeline", [
            "Filters fake registrations by mandating faculty upload valid institutional documents.",
            "Protects system integrity and user trust by locking account status until administrator clearance is complete."
        ], C_RED),
        ("Real-time PostgreSQL listeners", [
            "Integrates Supabase real-time triggers to synchronize scheduling status instantly.",
            "Avoids scheduling overlaps by immediately blocking and color-coding taken hour blocks."
        ], C_BLUE),
        ("Direct telemetry version checker", [
            "Eliminates client-side caching glitches by auto-checking version coordinates.",
            "Prompts active users to reload and refresh assets, avoiding dashboard session flickers."
        ], C_EMERALD)
    ]
    
    for idx, con in enumerate(contribs):
        x = Inches(0.5 + idx * 4.2)
        y = Inches(1.8)
        
        draw_rect(slide, x, y, Inches(3.9), Inches(5.0), C_NAVY_LIGHT, C_BLUE, 1)
        draw_rect(slide, x, y, Inches(3.9), Inches(0.12), con[2])
        
        tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.3), Inches(3.5), Inches(4.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = con[0]
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = C_MINT
        p.space_after = Pt(18)
        
        create_textbox_list(slide, x + Inches(0.2), y + Inches(1.0), Inches(3.5), Inches(3.7), con[1], font_size=12, font_color=C_WHITE, bold_title_prefix=False)

    # ----------------------------------------------------
    # SLIDE 15: CONCLUSION (Light Gray)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "CONCLUSION", "Summary of system design benefits and outcomes")
    
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(12.333), Inches(5.0), C_WHITE, C_GRAY_LIGHT)
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(0.15), Inches(5.0), C_NAVY_LIGHT)
    
    tb_concl = slide.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(11.3), Inches(4.4))
    tf_concl = tb_concl.text_frame
    tf_concl.word_wrap = True
    
    p = tf_concl.paragraphs[0]
    p.text = "Project Outcomes & System Integrity:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = C_NAVY_LIGHT
    p.space_after = Pt(14)
    
    concl_bullets = [
        "Eliminated Scheduling Conflicts: Displaying real-time calendars prevents coordination friction and double-bookings.",
        "Significant Efficiency Gains: Latency analysis confirms a **99.7% reduction** in scheduling search times compared to manual methods.",
        "Excellent User Usability: The System Usability Scale average score of **88.5 out of 100** verifies outstanding user satisfaction.",
        "Academic Trust Maintained: Admin vetting approval flow effectively locks out malicious or unauthorized accounts.",
        "Mobile Deployment Secured: Packaging client code via Capacitor wrapper allows smooth execution on Android devices."
    ]
    
    create_textbox_list(slide, Inches(1.0), Inches(2.7), Inches(11.3), Inches(3.8), concl_bullets, font_size=13.5)

    # ----------------------------------------------------
    # SLIDE 16: THANK YOU (Dark Navy)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_NAVY_DARK)
    
    # Background glowing accents
    draw_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), C_BLUE)
    draw_rect(slide, Inches(13.18), Inches(0), Inches(0.15), Inches(7.5), C_EMERALD)
    
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "THANK YOU!"
    p1.font.name = "Arial"
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = C_MINT
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(14)
    
    p2 = tf.add_paragraph()
    p2.text = "ConsulTime: Web & Mobile-Based Academic Consultation Scheduler"
    p2.font.name = "Arial"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_after = Pt(36)
    
    p3 = tf.add_paragraph()
    p3.text = "Questions & Answers Session"
    p3.font.name = "Arial"
    p3.font.size = Pt(16)
    p3.font.bold = True
    p3.font.color.rgb = C_BLUE
    p3.alignment = PP_ALIGN.CENTER
    
    p4 = tf.add_paragraph()
    p4.text = "Presented by: Khyle Bulawan, Carlvyn Bajala, Jake Pinggol, Henan Oliveros"
    p4.font.name = "Arial"
    p4.font.size = Pt(13)
    p4.font.color.rgb = C_TEXT_MUTED
    p4.alignment = PP_ALIGN.CENTER
    p4.space_before = Pt(20)

    # Save to the user's workspace target filename
    prs.save("ConsulTime PPT.pptx")
    print("SUCCESS: Programmatically compiled ConsulTime PPT.pptx slide deck!")

if __name__ == "__main__":
    create_presentation()
