import os
from pptx import Presentation
from pptx.util import Inches

def insert_screenshots():
    pptx_path = "ConsulTime_Modern_Presentation.pptx"
    if not os.path.exists(pptx_path):
        print(f"Error: '{pptx_path}' not found. Please run the styling script first.")
        return

    prs = Presentation(pptx_path)

    # Image mapping
    image_map = {
        6: "login_screenshot.png",       # Slide 6: Login Module
        7: "register_screenshot.png",    # Slide 7: Registration Module
        8: "student_screenshot.png",     # Slide 8: Student Dashboard
        9: "student_screenshot.png",     # Slide 9: Book Appointment (using Student Dashboard as fallback)
        10: "faculty_screenshot.png",    # Slide 10: Faculty Dashboard
        11: "admin_screenshot.png"       # Slide 11: Admin Dashboard
    }

    for idx, slide in enumerate(prs.slides):
        slide_num = idx + 1
        
        # Check if we have an image for this slide
        if slide_num in image_map:
            image_name = image_map[slide_num]
            if not os.path.exists(image_name):
                print(f"Warning: Image file '{image_name}' not found. Skipping slide {slide_num}.")
                continue

            # Look for placeholder text in shapes
            target_shape = None
            for shape in list(slide.shapes):
                if hasattr(shape, "text") and any(phrase in shape.text for phrase in ["Screenshot Here", "PLACE SCREENSHOT", "Screenshot"]):
                    target_shape = shape
                    break

            if target_shape:
                # Get position and size of the text placeholder
                left = target_shape.left
                top = target_shape.top
                width = target_shape.width
                height = target_shape.height

                # Delete the placeholder text box shape
                spTree = slide.shapes._spTree
                spTree.remove(target_shape._element)

                # Add the actual image in its place
                # We can also add a nice border frame around it
                slide.shapes.add_picture(image_name, left, top, width, height)
                print(f"SUCCESS: Inserted '{image_name}' into Slide {slide_num}")
            else:
                # If no placeholder text shape was found, let's find the first white rectangle card on the right
                # to insert the image inside it.
                # In Slide 6-11, the right side usually has a large card or space.
                # Let's add it at a default position if no placeholder is found
                left = Inches(5.4)
                top = Inches(1.8)
                width = Inches(7.4)
                height = Inches(5.0)
                slide.shapes.add_picture(image_name, left, top, width, height)
                print(f"SUCCESS: Placed '{image_name}' in default right-side area on Slide {slide_num}")

    # Save the modified presentation
    prs.save("ConsulTime_Modern_Presentation.pptx")
    print("SUCCESS: Finished inserting all screenshots into ConsulTime_Modern_Presentation.pptx!")

if __name__ == "__main__":
    insert_screenshots()
