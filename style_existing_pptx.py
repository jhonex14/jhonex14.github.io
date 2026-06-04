import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def style_presentation():
    target_file = "ConsulTime_Modern_Presentation.pptx"
    
    if not os.path.exists(target_file):
        print(f"Error: '{target_file}' not found.")
        sys.exit(1)
            
    prs = Presentation(target_file)
    # Widescreen proportions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette Constants
    C_NAVY_DARK  = RGBColor(15, 23, 42)      # #0F172A (Slate 900)
    C_NAVY_LIGHT = RGBColor(30, 58, 138)     # #1E3A8A (Navy)
    C_BLUE       = RGBColor(59, 130, 246)     # #3B82F6 (Blue)
    C_EMERALD    = RGBColor(16, 185, 129)     # #10B981 (Emerald)
    C_MINT       = RGBColor(52, 211, 153)     # #34D399 (Mint)
    C_BG_LIGHT   = RGBColor(248, 250, 252)    # #F8FAFC (Slate 50)
    C_WHITE      = RGBColor(255, 255, 255)    # #FFFFFF
    C_TEXT_DARK  = RGBColor(31, 41, 55)      # #1F2937 (Slate 800)
    C_TEXT_MUTED = RGBColor(107, 114, 128)    # #6B7280 (Gray)
    C_GRAY_LIGHT = RGBColor(229, 231, 235)    # #E5E7EB (Gray 200)

    # Z-Ordering helper: move a shape to the back of the slide shapes tree
    def send_to_back(slide, shape, pos=2):
        spTree = slide.shapes._spTree
        spTree.remove(shape._element)
        spTree.insert(pos, shape._element)

    dark_slide_indices = [0, 3, 13, 15] # Slides 1, 4, 14, 16 are dark navy themed

    for idx, slide in enumerate(prs.slides):
        is_dark = (idx in dark_slide_indices)
        
        # 1. Clear layout placeholders if they are blank/white shapes or backgrounds
        # Draw slide background solid rectangle covering the entire page
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_NAVY_DARK if is_dark else C_BG_LIGHT
        bg.line.fill.background()
        send_to_back(slide, bg, 2) # Position 2 is the bottom of shape tree

        # 2. Add decorative borders & accent bars
        if is_dark:
            # Left vertical emerald bar
            bar1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), Inches(7.5))
            bar1.fill.solid()
            bar1.fill.fore_color.rgb = C_EMERALD
            bar1.line.fill.background()
            send_to_back(slide, bar1, 3)
            
            # Right vertical blue bar
            bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(13.18), 0, Inches(0.15), Inches(7.5))
            bar2.fill.solid()
            bar2.fill.fore_color.rgb = C_BLUE
            bar2.line.fill.background()
            send_to_back(slide, bar2, 3)
        else:
            # Top horizontal Navy bar
            bar1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.08))
            bar1.fill.solid()
            bar1.fill.fore_color.rgb = C_NAVY_LIGHT
            bar1.line.fill.background()
            send_to_back(slide, bar1, 3)
            
            # Left vertical blue stripe
            stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.08), Inches(7.5))
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = C_BLUE
            stripe.line.fill.background()
            send_to_back(slide, stripe, 3)

        # Gather original shapes on the slide (excluding the ones we just added)
        new_shapes_list = [bg, bar1]
        if is_dark:
            new_shapes_list.append(bar2)
        else:
            new_shapes_list.append(stripe)
            
        original_shapes = [s for s in slide.shapes if s not in new_shapes_list]

        # 3. Format and design each shape
        card_position = 4 # Index in spTree to render cards above background but below text
        
        for shape in original_shapes:
            # Check if the shape has a text frame
            if hasattr(shape, "text_frame") and shape.text_frame:
                tf = shape.text_frame
                tf.word_wrap = True
                
                # Check if this shape is the slide's main title
                is_title = False
                if shape.top < Inches(1.6) and len(shape.text.strip()) < 100:
                    is_title = True
                
                # Format paragraphs
                for paragraph in tf.paragraphs:
                    paragraph.font.name = "Arial"
                    
                    if is_title:
                        paragraph.font.bold = True
                        paragraph.font.size = Pt(28) if not is_dark else Pt(36)
                        paragraph.font.color.rgb = C_MINT if is_dark else C_NAVY_LIGHT
                        
                        # Left align titles on light slides, Center on Title/Q&A
                        if idx in [0, 15]:
                            paragraph.alignment = PP_ALIGN.CENTER
                        else:
                            paragraph.alignment = PP_ALIGN.LEFT
                    else:
                        # Style body paragraphs and list bullets
                        if paragraph.font.size is None or paragraph.font.size < Pt(13):
                            paragraph.font.size = Pt(13)
                        paragraph.font.color.rgb = C_WHITE if is_dark else C_TEXT_DARK
                        
                        # Style individual text runs inside paragraph
                        for run in paragraph.runs:
                            run.font.name = "Arial"
                            run.font.color.rgb = C_WHITE if is_dark else C_TEXT_DARK
                
                # Style background containers for body text frames (make them into cards!)
                if not is_title and len(shape.text.strip()) > 3:
                    # Draw a rounded rectangle to act as a card container behind the text block
                    card = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        shape.left - Inches(0.15),
                        shape.top - Inches(0.15),
                        shape.width + Inches(0.3),
                        shape.height + Inches(0.3)
                    )
                    card.fill.solid()
                    if is_dark:
                        card.fill.fore_color.rgb = RGBColor(30, 41, 59) # Slate card
                        card.line.color.rgb = RGBColor(71, 85, 105)
                        card.line.width = Pt(1)
                    else:
                        card.fill.fore_color.rgb = C_WHITE # White card
                        card.line.color.rgb = C_GRAY_LIGHT
                        card.line.width = Pt(1.5)
                        
                    # Send card to Z-order position above background
                    send_to_back(slide, card, card_position)
                    card_position += 1

            # Style Table shapes (e.g. system evaluation results)
            elif shape.has_table:
                table = shape.table
                for r_idx, row in enumerate(table.rows):
                    for c_idx, cell in enumerate(row.cells):
                        cell.fill.solid()
                        if r_idx == 0:
                            # Header styling
                            cell.fill.fore_color.rgb = C_NAVY_LIGHT
                            for p in cell.text_frame.paragraphs:
                                p.font.name = "Arial"
                                p.font.bold = True
                                p.font.size = Pt(11)
                                p.font.color.rgb = C_WHITE
                                p.alignment = PP_ALIGN.CENTER
                        else:
                            # Data rows styling
                            cell.fill.fore_color.rgb = C_WHITE
                            for p in cell.text_frame.paragraphs:
                                p.font.name = "Arial"
                                p.font.size = Pt(10.5)
                                p.font.color.rgb = C_TEXT_DARK
                                p.alignment = PP_ALIGN.CENTER

            # Style Picture shapes (screenshots or logos)
            elif shape.shape_type == 13: # 13 is MSO_SHAPE_TYPE.PICTURE
                # Draw a clean border frame behind the picture to give it a modern shadow-frame mount look
                frame = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    shape.left - Inches(0.08),
                    shape.top - Inches(0.08),
                    shape.width + Inches(0.16),
                    shape.height + Inches(0.16)
                )
                frame.fill.solid()
                if is_dark:
                    frame.fill.fore_color.rgb = RGBColor(30, 41, 59)
                    frame.line.color.rgb = RGBColor(71, 85, 105)
                else:
                    frame.fill.fore_color.rgb = C_WHITE
                    frame.line.color.rgb = C_GRAY_LIGHT
                    frame.line.width = Pt(1.5)
                    
                send_to_back(slide, frame, card_position)
                card_position += 1

    output_file = "ConsulTime_Modern_Presentation.pptx"
    prs.save(output_file)
    print(f"SUCCESS: Successfully styled existing slides in-place and saved to '{output_file}'!")

if __name__ == "__main__":
    style_presentation()
