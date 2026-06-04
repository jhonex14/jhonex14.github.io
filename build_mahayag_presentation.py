import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ----------------------------------------------------
    # COLOR PALETTE SETUP
    # ----------------------------------------------------
    C_NAVY_DARK  = RGBColor(15, 23, 42)      # #0F172A
    C_NAVY_LIGHT = RGBColor(30, 58, 138)     # #1E3A8A
    C_BLUE       = RGBColor(59, 130, 246)     # #3B82F6
    C_EMERALD    = RGBColor(16, 185, 129)     # #10B981
    C_MINT       = RGBColor(52, 211, 153)     # #34D399
    C_BG_LIGHT   = RGBColor(248, 250, 252)    # #F8FAFC
    C_WHITE      = RGBColor(255, 255, 255)    # #FFFFFF
    C_TEXT_DARK  = RGBColor(31, 41, 55)      # #1F2937
    C_TEXT_MUTED = RGBColor(107, 114, 128)    # #6B7280
    C_GRAY_LIGHT = RGBColor(229, 231, 235)    # #E5E7EB
    C_RED        = RGBColor(239, 68, 68)      # #EF4444

    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # Helper function to add slide background
    def set_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper to draw colored rectangle shapes (cards, panels, accents)
    def draw_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=1):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(border_width)
        else:
            shape.line.fill.background()
        return shape

    # Helper to add modern slide titles with primary & accent text styling
    def add_slide_header(slide, title_text, subtitle_text=None, is_dark_bg=False):
        # Draw top accent bar
        draw_rect(slide, Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.08), C_EMERALD if is_dark_bg else C_NAVY_LIGHT)
        
        # Add Header Textbox
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = C_WHITE if is_dark_bg else C_NAVY_LIGHT
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = "Arial"
            p2.font.size = Pt(14)
            p2.font.color.rgb = C_MINT if is_dark_bg else C_TEXT_MUTED

    # Helper to create styled bullet lists inside textboxes
    def create_textbox_list(slide, left, top, width, height, bullets, font_size=13, font_color=C_TEXT_DARK, bold_title_prefix=True):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        for idx, bullet in enumerate(bullets):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.font.name = "Arial"
            p.font.size = Pt(font_size)
            p.space_after = Pt(8)
            p.level = 0
            
            # Highlight bold prefix (e.g. "Main Task: explanation")
            if bold_title_prefix and ":" in bullet:
                parts = bullet.split(":", 1)
                run1 = p.add_run()
                run1.text = "• " + parts[0] + ":"
                run1.font.bold = True
                run1.font.color.rgb = font_color
                
                run2 = p.add_run()
                run2.text = parts[1]
                run2.font.bold = False
                run2.font.color.rgb = font_color
            else:
                p.text = "• " + bullet
                p.font.color.rgb = font_color
        return tb

    # Helper to insert custom screenshots with modern cards/borders
    def insert_screenshot(slide, image_name, left, top, width, height, is_dark_bg=False):
        if not os.path.exists(image_name):
            print(f"Warning: image '{image_name}' not found. Adding visual placeholder instead.")
            # Draw placeholder card
            draw_rect(slide, left, top, width, height, C_NAVY_LIGHT if is_dark_bg else C_WHITE, C_GRAY_LIGHT, 2)
            tb = slide.shapes.add_textbox(left, top + height/2 - Inches(0.5), width, Inches(1))
            p = tb.text_frame.paragraphs[0]
            p.text = f"Screenshot Placeholder:\n{image_name}"
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Arial"
            p.font.size = Pt(14)
            p.font.color.rgb = C_WHITE if is_dark_bg else C_TEXT_DARK
        else:
            # Draw clean border frame behind the picture to give it a modern shadow-frame mount look
            frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left - Inches(0.08), top - Inches(0.08), width + Inches(0.16), height + Inches(0.16))
            frame.fill.solid()
            if is_dark_bg:
                frame.fill.fore_color.rgb = RGBColor(30, 41, 59)
                frame.line.color.rgb = RGBColor(71, 85, 105)
            else:
                frame.fill.fore_color.rgb = C_WHITE
                frame.line.color.rgb = C_GRAY_LIGHT
                frame.line.width = Pt(1.5)
            
            # Add the actual image
            slide.shapes.add_picture(image_name, left, top, width, height)

    # ----------------------------------------------------
    # SLIDE 1: TITLE SLIDE (Dark Navy)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_NAVY_DARK)
    
    # Background glowing accents
    draw_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), C_EMERALD)
    draw_rect(slide, Inches(13.18), Inches(0), Inches(0.15), Inches(7.5), C_BLUE)
    
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "CONSULTIME"
    p1.font.name = "Arial"
    p1.font.size = Pt(56)
    p1.font.bold = True
    p1.font.color.rgb = C_MINT
    p1.space_after = Pt(12)
    
    p2 = tf.add_paragraph()
    p2.text = "A Cross-Platform Consultation and Appointment Management System for Mahayag National High School"
    p2.font.name = "Arial"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE
    p2.space_after = Pt(40)
    
    p3 = tf.add_paragraph()
    p3.text = "Presented by: Philipp Edward Sapalicio, Ednell Sapalicio"
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = C_WHITE
    
    p4 = tf.add_paragraph()
    p4.text = "ITE Department, Philippine College of Technology  |  May 2026"
    p4.font.name = "Arial"
    p4.font.size = Pt(13)
    p4.font.color.rgb = C_BLUE

    # ----------------------------------------------------
    # SLIDE 2: PROBLEM STATEMENT (Light Gray)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "PROBLEM STATEMENT", "Challenges in manual scheduling and communication at Mahayag NHS")
    
    # Left Column: Context Card
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(4.0), Inches(5.0), C_NAVY_DARK)
    tb_left = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(3.4), Inches(4.2))
    tf_left = tb_left.text_frame
    tf_left.word_wrap = True
    p_left1 = tf_left.paragraphs[0]
    p_left1.text = "The Scheduling friction"
    p_left1.font.bold = True
    p_left1.font.size = Pt(18)
    p_left1.font.color.rgb = C_MINT
    p_left1.space_after = Pt(14)
    p_left2 = tf_left.add_paragraph()
    p_left2.text = "Academic scheduling at Mahayag National High School remains largely manual, relying on paper logbooks, physical bulletin boards, and uncoordinated verbal updates.\n\nThis leads to inefficient time management, missed consultation slots, and coordination gaps."
    p_left2.font.size = Pt(13)
    p_left2.font.color.rgb = C_WHITE
    p_left2.space_before = Pt(10)
    
    # Right Column: Four Key Pain Points
    problems = [
        "Uncoordinated Timetables: Relying on static lists on doors leads to appointment conflicts.",
        "Communication Gaps: Lack of instant notifications for sudden timetable changes.",
        "Security Concerns: No centralized verification process, allowing unvetted sign-ups.",
        "No Touch Optimization: Current systems lack mobile responsiveness, limiting accessibility."
    ]
    for idx, prob in enumerate(problems):
        row = idx // 2
        col = idx % 2
        x = Inches(4.85 + col * 4.0)
        y = Inches(1.8 + row * 2.5)
        
        draw_rect(slide, x, y, Inches(3.7), Inches(2.3), C_WHITE, C_GRAY_LIGHT)
        draw_rect(slide, x, y, Inches(0.08), Inches(2.3), C_RED)
        
        create_textbox_list(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.3), Inches(1.9), [prob], font_size=12.5)

    # ----------------------------------------------------
    # SLIDE 3: OBJECTIVES OF THE STUDY (Light Gray)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "OBJECTIVES", "Primary and specific objectives of ConsulTime")
    
    # Left Column: Primary Objective
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(4.5), Inches(5.0), C_NAVY_LIGHT)
    tb_obj = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(3.9), Inches(4.4))
    tf_obj = tb_obj.text_frame
    tf_obj.word_wrap = True
    p_obj1 = tf_obj.paragraphs[0]
    p_obj1.text = "Primary Objective"
    p_obj1.font.bold = True
    p_obj1.font.size = Pt(20)
    p_obj1.font.color.rgb = C_MINT
    p_obj1.space_after = Pt(20)
    p_obj2 = tf_obj.add_paragraph()
    p_obj2.text = "To design, develop, and implement ConsulTime: a real-time academic consultation and appointment management system for Mahayag National High School to digitize and streamline scheduling workflows."
    p_obj2.font.size = Pt(14)
    p_obj2.font.color.rgb = C_WHITE
    p_obj2.space_before = Pt(10)
    
    # Right Column: Specific Objectives
    draw_rect(slide, Inches(5.3), Inches(1.8), Inches(7.5), Inches(5.0), C_WHITE, C_GRAY_LIGHT)
    draw_rect(slide, Inches(5.3), Inches(1.8), Inches(0.12), Inches(5.0), C_EMERALD)
    
    tb_spec = slide.shapes.add_textbox(Inches(5.7), Inches(2.0), Inches(6.8), Inches(4.6))
    tf_spec = tb_spec.text_frame
    tf_spec.word_wrap = True
    p_spec_title = tf_spec.paragraphs[0]
    p_spec_title.text = "Specific Capstone Targets:"
    p_spec_title.font.bold = True
    p_spec_title.font.size = Pt(17)
    p_spec_title.font.color.rgb = C_NAVY_LIGHT
    p_spec_title.space_after = Pt(12)
    
    spec_bullets = [
        "Secure User Gateways: Develop secure registration and login views for Students, Faculty, and Administrators.",
        "Dynamic Schedule Blocks: Allow faculty members to declare, modify, and close availability hours in real time.",
        "Responsive Client UI: Establish a mobile-friendly appointment booking form for students.",
        "Account Vetting Panel: Integrate approval toggles for admins to verify faculty credentials before activation.",
        "Cross-Platform Compilation: Deploy system as a local-cached PWA and native Android App wrapper (.APK)."
    ]
    create_textbox_list(slide, Inches(5.7), Inches(2.5), Inches(6.8), Inches(4.1), spec_bullets, font_size=12.5)

    # ----------------------------------------------------
    # SLIDE 4: SCOPE AND LIMITATIONS (Light Gray)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "SCOPE AND LIMITATIONS", "Technical boundaries and operational constraints of the system")
    
    # Left Column: Scope
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), C_WHITE, C_GRAY_LIGHT)
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(6.0), Inches(0.12), C_BLUE)
    
    tb_scope = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(5.4), Inches(4.4))
    tf_scope = tb_scope.text_frame
    tf_scope.word_wrap = True
    p_sc = tf_scope.paragraphs[0]
    p_sc.text = "Project Scope"
    p_sc.font.bold = True
    p_sc.font.size = Pt(16)
    p_sc.font.color.rgb = C_NAVY_LIGHT
    p_sc.space_after = Pt(14)
    
    scope_bullets = [
        "Scheduling Engine: Tracks student appointment requests and faculty responses.",
        "Real-Time Sync: Updates calendar states instantly across client dashboards when changed.",
        "Authentication Controls: Limits system read/write capabilities according to user roles.",
        "Administrative Vetting: Validates faculty profile registration parameters."
    ]
    create_textbox_list(slide, Inches(0.8), Inches(2.7), Inches(5.4), Inches(3.7), scope_bullets, font_size=12.5)
    
    # Right Column: Limitations
    draw_rect(slide, Inches(6.8), Inches(1.8), Inches(6.0), Inches(5.0), C_WHITE, C_GRAY_LIGHT)
    draw_rect(slide, Inches(6.8), Inches(1.8), Inches(6.0), Inches(0.12), C_RED)
    
    tb_lim = slide.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.4), Inches(4.4))
    tf_lim = tb_lim.text_frame
    tf_lim.word_wrap = True
    p_lm = tf_lim.paragraphs[0]
    p_lm.text = "System Limitations"
    p_lm.font.bold = True
    p_lm.font.size = Pt(16)
    p_lm.font.color.rgb = C_NAVY_LIGHT
    p_lm.space_after = Pt(14)
    
    lim_bullets = [
        "No Video Calls: Stakeholders communicate inside the system but meet physically (no video frame integration).",
        "iOS Limitations: Mobile native compilation generates Android wrapper packages (.APK) only.",
        "Internet Dependency: Relies on cloud listeners and requires active connections for data sync.",
        "No Document Scanning: Relies on user uploads (no optical character recognition utilities)."
    ]
    create_textbox_list(slide, Inches(7.1), Inches(2.7), Inches(5.4), Inches(3.7), lim_bullets, font_size=12.5)

    # ----------------------------------------------------
    # SLIDE 5: SIGNIFICANCE OF THE STUDY (Light Gray)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "SIGNIFICANCE OF THE STUDY", "Impact of ConsulTime on school stakeholders")
    
    significances = [
        ("To Students", [
            "Reduces manual scheduling friction.",
            "Enables quick reservation of slots directly from phones.",
            "Keeps students updated on request statuses instantly."
        ], C_BLUE),
        ("To Faculty Members", [
            "Saves time by removing manual logbooks.",
            "Allows setting custom consultation availability slots.",
            "Prevents double-booking conflicts automatically."
        ], C_EMERALD),
        ("To School Administrators", [
            "Provides unified portal to audit registrations.",
            "Maintains verified, secure list of active academic roles.",
            "Saves history of consultations for research records."
        ], C_NAVY_LIGHT)
    ]
    
    for idx, sig in enumerate(significances):
        x = Inches(0.5 + idx * 4.2)
        y = Inches(1.8)
        
        draw_rect(slide, x, y, Inches(3.9), Inches(5.0), C_WHITE, C_GRAY_LIGHT)
        draw_rect(slide, x, y, Inches(3.9), Inches(0.12), sig[2])
        
        tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.3), Inches(3.5), Inches(4.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = sig[0]
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = C_NAVY_LIGHT
        p.space_after = Pt(18)
        
        create_textbox_list(slide, x + Inches(0.2), y + Inches(1.0), Inches(3.5), Inches(3.7), sig[1], font_size=12, bold_title_prefix=False)

    # ----------------------------------------------------
    # SLIDE 6: SYSTEM ARCHITECTURE (Dark Navy)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_NAVY_DARK)
    add_slide_header(slide, "SYSTEM ARCHITECTURE", "ConsulTime's three-tier design layout", is_dark_bg=True)
    
    # Client Layer
    draw_rect(slide, Inches(0.8), Inches(2.5), Inches(3.2), Inches(3.5), C_NAVY_LIGHT, C_BLUE, 2)
    tb_arch1 = slide.shapes.add_textbox(Inches(1.0), Inches(2.7), Inches(2.8), Inches(3.1))
    tf_arch1 = tb_arch1.text_frame
    tf_arch1.word_wrap = True
    tf_arch1.paragraphs[0].text = "CLIENT TIER"
    tf_arch1.paragraphs[0].font.bold = True
    tf_arch1.paragraphs[0].font.size = Pt(15)
    tf_arch1.paragraphs[0].font.color.rgb = C_MINT
    tf_arch1.paragraphs[0].space_after = Pt(10)
    create_textbox_list(slide, Inches(0.9), Inches(3.2), Inches(3.0), Inches(2.7), [
        "Web App: HTML5, CSS3, ES6 JS.",
        "Android Package: Compiled native wrapper via Capacitor.",
        "Caching: Service Worker (sw.js) for fast assets offline reload."
    ], font_size=11.5, font_color=C_WHITE, bold_title_prefix=False)

    # Arrow 1
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.2), Inches(4.0), Inches(0.7), Inches(0.5))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = C_BLUE
    arrow1.line.fill.background()

    # Backend Layer
    draw_rect(slide, Inches(5.1), Inches(2.5), Inches(3.2), Inches(3.5), C_NAVY_LIGHT, C_EMERALD, 2)
    tb_arch2 = slide.shapes.add_textbox(Inches(5.3), Inches(2.7), Inches(2.8), Inches(3.1))
    tf_arch2 = tb_arch2.text_frame
    tf_arch2.word_wrap = True
    tf_arch2.paragraphs[0].text = "CLOUD DATABASE TIER"
    tf_arch2.paragraphs[0].font.bold = True
    tf_arch2.paragraphs[0].font.size = Pt(15)
    tf_arch2.paragraphs[0].font.color.rgb = C_MINT
    tf_arch2.paragraphs[0].space_after = Pt(10)
    create_textbox_list(slide, Inches(5.2), Inches(3.2), Inches(3.0), Inches(2.7), [
        "Supabase Instance: Relational PostgreSQL schema host.",
        "Row Level Security (RLS): Enforces access restrictions.",
        "Real-Time Subscriptions: Pushes database writes instantly."
    ], font_size=11.5, font_color=C_WHITE, bold_title_prefix=False)

    # Arrow 2
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.5), Inches(4.0), Inches(0.7), Inches(0.5))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = C_EMERALD
    arrow2.line.fill.background()

    # Admin Control Layer
    draw_rect(slide, Inches(9.4), Inches(2.5), Inches(3.2), Inches(3.5), C_NAVY_LIGHT, C_BLUE, 2)
    tb_arch3 = slide.shapes.add_textbox(Inches(9.6), Inches(2.7), Inches(2.8), Inches(3.1))
    tf_arch3 = tb_arch3.text_frame
    tf_arch3.word_wrap = True
    tf_arch3.paragraphs[0].text = "ADMINISTRATIVE TIER"
    tf_arch3.paragraphs[0].font.bold = True
    tf_arch3.paragraphs[0].font.size = Pt(15)
    tf_arch3.paragraphs[0].font.color.rgb = C_MINT
    tf_arch3.paragraphs[0].space_after = Pt(10)
    create_textbox_list(slide, Inches(9.5), Inches(3.2), Inches(3.0), Inches(2.7), [
        "Vetting Controls: Activates pending faculty registrations.",
        "Telemetry Checks: Monitors database speeds and latencies.",
        "Version Control: Syncs version.json to reload clients."
    ], font_size=11.5, font_color=C_WHITE, bold_title_prefix=False)

    # ----------------------------------------------------
    # SLIDE 7: DATABASE SCHEMA I (Light Gray)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "DATABASE SCHEMA: USERS & AVAILABILITIES", "Detailed view of the profiles and availabilities relations")
    
    # Left Card: User Profiles
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), C_WHITE, C_GRAY_LIGHT)
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(6.0), Inches(0.12), C_NAVY_LIGHT)
    
    tb_p1 = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(5.4), Inches(4.4))
    tf_p1 = tb_p1.text_frame
    tf_p1.word_wrap = True
    tf_p1.paragraphs[0].text = "Table 1: User Profiles (profiles)"
    tf_p1.paragraphs[0].font.bold = True
    tf_p1.paragraphs[0].font.size = Pt(16)
    tf_p1.paragraphs[0].font.color.rgb = C_NAVY_LIGHT
    tf_p1.paragraphs[0].space_after = Pt(10)
    
    profiles_bullets = [
        "uuid id: Primary Key linked to Supabase auth credentials.",
        "text role: Limits users to Student, Faculty, or Admin.",
        "text department: Stores academic department information.",
        "text id_number: University standard identification field.",
        "boolean is_approved: Vetting flag for faculty accounts."
    ]
    create_textbox_list(slide, Inches(0.8), Inches(2.7), Inches(5.4), Inches(3.7), profiles_bullets, font_size=12.5)

    # Right Card: Faculty Availabilities
    draw_rect(slide, Inches(6.8), Inches(1.8), Inches(6.0), Inches(5.0), C_WHITE, C_GRAY_LIGHT)
    draw_rect(slide, Inches(6.8), Inches(1.8), Inches(6.0), Inches(0.12), C_BLUE)
    
    tb_p2 = slide.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.4), Inches(4.4))
    tf_p2 = tb_p2.text_frame
    tf_p2.word_wrap = True
    tf_p2.paragraphs[0].text = "Table 2: Availabilities (availabilities)"
    tf_p2.paragraphs[0].font.bold = True
    tf_p2.paragraphs[0].font.size = Pt(16)
    tf_p2.paragraphs[0].font.color.rgb = C_NAVY_LIGHT
    tf_p2.paragraphs[0].space_after = Pt(10)
    
    avail_bullets = [
        "uuid faculty_id: Foreign Key pointing to profiles.id.",
        "int4 day_of_week: Maps recurring days from Monday to Friday.",
        "time start_time: Availability block start boundary.",
        "time end_time: Availability block end boundary.",
        "date specific_date: Optional single-occurrence time window overrides."
    ]
    create_textbox_list(slide, Inches(7.1), Inches(2.7), Inches(5.4), Inches(3.7), avail_bullets, font_size=12.5)

    # ----------------------------------------------------
    # SLIDE 8: DATABASE SCHEMA II (Light Gray)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "DATABASE SCHEMA: BOOKINGS", "Transactional engine of the system")
    
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(12.333), Inches(5.0), C_WHITE, C_GRAY_LIGHT)
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(0.15), Inches(5.0), C_EMERALD)
    
    tb_bk = slide.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(11.3), Inches(4.4))
    tf_bk = tb_bk.text_frame
    tf_bk.word_wrap = True
    tf_bk.paragraphs[0].text = "Table 3: Consultation Bookings (bookings)"
    tf_bk.paragraphs[0].font.bold = True
    tf_bk.paragraphs[0].font.size = Pt(18)
    tf_bk.paragraphs[0].font.color.rgb = C_NAVY_LIGHT
    tf_bk.paragraphs[0].space_after = Pt(14)
    
    bookings_bullets = [
        "uuid student_id & faculty_id: Dual Foreign Keys referencing user profile tables to bind participants.",
        "date appointment_date: Calendar date of the consultation.",
        "time start_time & end_time: Precise timing parameters of the consultation.",
        "text purpose: Qualitative reason block filled out by the student.",
        "text status: Booking state constraint strictly restricted to Pending, Approved, or Rejected.",
        "text faculty_notes: Optional feedback block for instructors to write resolutions."
    ]
    create_textbox_list(slide, Inches(1.0), Inches(2.7), Inches(11.3), Inches(3.8), bookings_bullets, font_size=13.5)

    # ----------------------------------------------------
    # SLIDE 9: DATABASE RELATIONAL TABLES (Light Gray)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "DATABASE RELATIONAL TABLES", "Live database table records hosted on Supabase Cloud")
    
    # Left screenshot: Profiles Table
    insert_screenshot(slide, "db_profiles_screenshot.png", Inches(0.5), Inches(1.8), Inches(3.8), Inches(3.6))
    tb_db1 = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(3.8), Inches(1.7))
    tf_db1 = tb_db1.text_frame
    tf_db1.word_wrap = True
    p_db1 = tf_db1.paragraphs[0]
    p_db1.text = "Supabase profiles schema showing role distribution, credentials mapping, and user vetting statuses."
    p_db1.font.name = "Arial"
    p_db1.font.size = Pt(11)
    p_db1.font.color.rgb = C_TEXT_MUTED
    p_db1.alignment = PP_ALIGN.CENTER
    
    # Center screenshot: Availabilities Table
    insert_screenshot(slide, "db_availability_screenshot.png", Inches(4.76), Inches(1.8), Inches(3.8), Inches(3.6))
    tb_db2 = slide.shapes.add_textbox(Inches(4.76), Inches(5.5), Inches(3.8), Inches(1.7))
    tf_db2 = tb_db2.text_frame
    tf_db2.word_wrap = True
    p_db2 = tf_db2.paragraphs[0]
    p_db2.text = "Faculty availability slots table with dynamic timestamps, day-of-week markers, and recurring boundaries."
    p_db2.font.name = "Arial"
    p_db2.font.size = Pt(11)
    p_db2.font.color.rgb = C_TEXT_MUTED
    p_db2.alignment = PP_ALIGN.CENTER
    
    # Right screenshot: Bookings Table
    insert_screenshot(slide, "db_bookings_screenshot.png", Inches(9.02), Inches(1.8), Inches(3.8), Inches(3.6))
    tb_db3 = slide.shapes.add_textbox(Inches(9.02), Inches(5.5), Inches(3.8), Inches(1.7))
    tf_db3 = tb_db3.text_frame
    tf_db3.word_wrap = True
    p_db3 = tf_db3.paragraphs[0]
    p_db3.text = "Transactional bookings junction table mapping student/instructor requests and status state transitions."
    p_db3.font.name = "Arial"
    p_db3.font.size = Pt(11)
    p_db3.font.color.rgb = C_TEXT_MUTED
    p_db3.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # SLIDE 10: SHOWCASE: ACCESS & REGISTRATION (Light Gray)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "SYSTEM SHOWCASE: ACCESS & REGISTRATION", "Secure user gateway and sign-up pipeline views")
    
    # Left screenshot: Sign In
    insert_screenshot(slide, "login_screenshot.png", Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.5))
    tb_lbl1 = slide.shapes.add_textbox(Inches(1.0), Inches(6.4), Inches(5.2), Inches(0.8))
    p_lbl1 = tb_lbl1.text_frame.paragraphs[0]
    p_lbl1.text = "Figure 1: Login interface. Routes user role to Student, Faculty, or Admin dashboard automatically."
    p_lbl1.font.name = "Arial"
    p_lbl1.font.size = Pt(11)
    p_lbl1.font.color.rgb = C_TEXT_MUTED
    p_lbl1.alignment = PP_ALIGN.CENTER
    
    # Right screenshot: Register
    insert_screenshot(slide, "register_screenshot.png", Inches(7.1), Inches(1.8), Inches(5.2), Inches(4.5))
    tb_lbl2 = slide.shapes.add_textbox(Inches(7.1), Inches(6.4), Inches(5.2), Inches(0.8))
    p_lbl2 = tb_lbl2.text_frame.paragraphs[0]
    p_lbl2.text = "Figure 2: Registration console. Restricts emails to school domains and collects faculty verification uploads."
    p_lbl2.font.name = "Arial"
    p_lbl2.font.size = Pt(11)
    p_lbl2.font.color.rgb = C_TEXT_MUTED
    p_lbl2.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # SLIDE 11: SHOWCASE: STUDENT DASHBOARD (Light Gray)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "SYSTEM SHOWCASE: STUDENT DASHBOARD", "Student view for scheduling and app access")
    
    # Left Description
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(4.5), Inches(5.0), C_WHITE, C_GRAY_LIGHT)
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(4.5), Inches(0.12), C_BLUE)
    
    tb_std_desc = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(3.9), Inches(4.3))
    tf_std_desc = tb_std_desc.text_frame
    tf_std_desc.word_wrap = True
    tf_std_desc.paragraphs[0].text = "Student Panel Layout"
    tf_std_desc.paragraphs[0].font.bold = True
    tf_std_desc.paragraphs[0].font.size = Pt(18)
    tf_std_desc.paragraphs[0].font.color.rgb = C_NAVY_LIGHT
    tf_std_desc.paragraphs[0].space_after = Pt(14)
    
    std_bullets = [
        "Upcoming Bookings: Lists active appointment schedules and statuses.",
        "APK Download: Banners link directly to Capacitor mobile packages.",
        "Alerts Panel: Shows dynamic responses from instructors."
    ]
    create_textbox_list(slide, Inches(0.8), Inches(2.7), Inches(3.9), Inches(3.8), std_bullets, font_size=12.5)
    
    # Right Image
    insert_screenshot(slide, "student_screenshot.png", Inches(5.4), Inches(1.8), Inches(7.4), Inches(5.0))

    # ----------------------------------------------------
    # SLIDE 12: SHOWCASE: FACULTY & ADMIN (Light Gray)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "SYSTEM SHOWCASE: FACULTY & ADMIN PORTALS", "Governance panels for faculty and admins")
    
    # Left screenshot: Faculty
    insert_screenshot(slide, "faculty_screenshot.png", Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.5))
    tb_f_lbl = slide.shapes.add_textbox(Inches(1.0), Inches(6.4), Inches(5.2), Inches(0.8))
    p_f_lbl = tb_f_lbl.text_frame.paragraphs[0]
    p_f_lbl.text = "Figure 3: Faculty Dashboard. Allows instructors to toggle slots status and approve/reject bookings."
    p_f_lbl.font.name = "Arial"
    p_f_lbl.font.size = Pt(11)
    p_f_lbl.font.color.rgb = C_TEXT_MUTED
    p_f_lbl.alignment = PP_ALIGN.CENTER
    
    # Right screenshot: Admin
    insert_screenshot(slide, "admin_screenshot.png", Inches(7.1), Inches(1.8), Inches(5.2), Inches(4.5))
    tb_a_lbl = slide.shapes.add_textbox(Inches(7.1), Inches(6.4), Inches(5.2), Inches(0.8))
    p_a_lbl = tb_a_lbl.text_frame.paragraphs[0]
    p_a_lbl.text = "Figure 4: Admin Portal. Unified gateway to approve faculty registrations and monitor server logs."
    p_a_lbl.font.name = "Arial"
    p_a_lbl.font.size = Pt(11)
    p_a_lbl.font.color.rgb = C_TEXT_MUTED
    p_a_lbl.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # SLIDE 13: ADMIN CONTROLS & VETTING PIPELINE (Light Gray)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "ADMIN CONTROLS & VETTING PIPELINE", "Vetting interfaces to verify faculty credentials and monitor server status")
    
    # Left screenshot: Admin Overview
    insert_screenshot(slide, "admin_overview_screenshot.png", Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.0))
    tb_adv1 = slide.shapes.add_textbox(Inches(0.6), Inches(5.9), Inches(5.8), Inches(1.0))
    tf_adv1 = tb_adv1.text_frame
    tf_adv1.word_wrap = True
    p_adv1 = tf_adv1.paragraphs[0]
    p_adv1.text = "Figure 3.1: Admin Dashboard. Shows user statistics and quick actions to manage system state."
    p_adv1.font.name = "Arial"
    p_adv1.font.size = Pt(11)
    p_adv1.font.color.rgb = C_TEXT_MUTED
    p_adv1.alignment = PP_ALIGN.CENTER
    
    # Right screenshot: Admin Vetting
    insert_screenshot(slide, "admin_vetting_screenshot.png", Inches(6.9), Inches(1.8), Inches(5.8), Inches(4.0))
    tb_adv2 = slide.shapes.add_textbox(Inches(6.9), Inches(5.9), Inches(5.8), Inches(1.0))
    tf_adv2 = tb_adv2.text_frame
    tf_adv2.word_wrap = True
    p_adv2 = tf_adv2.paragraphs[0]
    p_adv2.text = "Figure 3.2: Faculty Approvals. Admin interface to approve or reject pending faculty registrations."
    p_adv2.font.name = "Arial"
    p_adv2.font.size = Pt(11)
    p_adv2.font.color.rgb = C_TEXT_MUTED
    p_adv2.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # SLIDE 14: APK MOBILE PACKAGING & DIRECT DOWNLOAD (Light Gray)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "APK MOBILE PACKAGING & DIRECT DOWNLOAD", "Capacitor hybrid native wrapper and file server deployment")
    
    # Left Card: Bullets discussion
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(4.0), Inches(5.0), C_WHITE, C_GRAY_LIGHT)
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(4.0), Inches(0.12), C_EMERALD)
    
    tb_apk_disc = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(3.6), Inches(4.6))
    tf_apk_disc = tb_apk_disc.text_frame
    tf_apk_disc.word_wrap = True
    p_apk_title = tf_apk_disc.paragraphs[0]
    p_apk_title.text = "Why the Mobile App Matters:"
    p_apk_title.font.bold = True
    p_apk_title.font.size = Pt(15)
    p_apk_title.font.color.rgb = C_NAVY_LIGHT
    p_apk_title.space_after = Pt(12)
    
    apk_disc_bullets = [
        "Capacitor Runtime: Compiles web codebase into lightweight native Android APK wrapper package.",
        "Offline Caching: Service workers preserve layout assets to operate without high mobile bandwidth.",
        "Bypass App Stores: Direct download button bypasses verification delays for quick campus deployment.",
        "Budget Device Friendly: 3.69MB package size ensures compatibility with low-end devices (e.g. Vivo Y17s)."
    ]
    create_textbox_list(slide, Inches(0.7), Inches(2.5), Inches(3.6), Inches(4.1), apk_disc_bullets, font_size=11.5)
    
    # Center screenshot: Download Button
    insert_screenshot(slide, "apk_download_button.png", Inches(4.8), Inches(1.8), Inches(3.8), Inches(3.2))
    tb_apk_lbl1 = slide.shapes.add_textbox(Inches(4.8), Inches(5.2), Inches(3.8), Inches(1.4))
    tf_apk_lbl1 = tb_apk_lbl1.text_frame
    tf_apk_lbl1.word_wrap = True
    p_apk_lbl1 = tf_apk_lbl1.paragraphs[0]
    p_apk_lbl1.text = "Figure 5.1: Direct Android App (.APK) download action button integrated on Login and Register forms."
    p_apk_lbl1.font.name = "Arial"
    p_apk_lbl1.font.size = Pt(11)
    p_apk_lbl1.font.color.rgb = C_TEXT_MUTED
    p_apk_lbl1.alignment = PP_ALIGN.CENTER
    
    # Right screenshot: MediaFire Page
    insert_screenshot(slide, "apk_mediafire_download.png", Inches(8.9), Inches(1.8), Inches(3.8), Inches(3.2))
    tb_apk_lbl2 = slide.shapes.add_textbox(Inches(8.9), Inches(5.2), Inches(3.8), Inches(1.4))
    tf_apk_lbl2 = tb_apk_lbl2.text_frame
    tf_apk_lbl2.word_wrap = True
    p_apk_lbl2 = tf_apk_lbl2.paragraphs[0]
    p_apk_lbl2.text = "Figure 5.2: Hosted MediaFire server page distributing compiled direct application package (consultime.apk)."
    p_apk_lbl2.font.name = "Arial"
    p_apk_lbl2.font.size = Pt(11)
    p_apk_lbl2.font.color.rgb = C_TEXT_MUTED
    p_apk_lbl2.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # SLIDE 15: SYSTEM EVALUATION METHODOLOGY (Light Gray)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "SYSTEM EVALUATION METHODOLOGY", "Usability testing framework and focus groups")
    
    # Left: Evaluation parameters
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), C_WHITE, C_GRAY_LIGHT)
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(6.0), Inches(0.12), C_NAVY_LIGHT)
    
    tb_meth1 = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(5.4), Inches(4.4))
    tf_meth1 = tb_meth1.text_frame
    tf_meth1.word_wrap = True
    tf_meth1.paragraphs[0].text = "Testing Standard: ISO/IEC 25010"
    tf_meth1.paragraphs[0].font.bold = True
    tf_meth1.paragraphs[0].font.size = Pt(16)
    tf_meth1.paragraphs[0].font.color.rgb = C_NAVY_LIGHT
    tf_meth1.paragraphs[0].space_after = Pt(14)
    
    meth1_bullets = [
        "Quality Characteristics: Functionality, Usability, Reliability, Performance Efficiency, Security.",
        "Survey Tool: Five-point Likert scale questionnaire.",
        "Analysis: Overall mean and quality interpretations computed."
    ]
    create_textbox_list(slide, Inches(0.8), Inches(2.7), Inches(5.4), Inches(3.7), meth1_bullets, font_size=13)

    # Right: Respondents count
    draw_rect(slide, Inches(6.8), Inches(1.8), Inches(6.0), Inches(5.0), C_WHITE, C_GRAY_LIGHT)
    draw_rect(slide, Inches(6.8), Inches(1.8), Inches(6.0), Inches(0.12), C_EMERALD)
    
    tb_meth2 = slide.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.4), Inches(4.4))
    tf_meth2 = tb_meth2.text_frame
    tf_meth2.word_wrap = True
    tf_meth2.paragraphs[0].text = "Focus Group Respondents (N = 26)"
    tf_meth2.paragraphs[0].font.bold = True
    tf_meth2.paragraphs[0].font.size = Pt(16)
    tf_meth2.paragraphs[0].font.color.rgb = C_NAVY_LIGHT
    tf_meth2.paragraphs[0].space_after = Pt(14)
    
    meth2_bullets = [
        "Students: 15 active users tested booking and notification flows.",
        "Faculty Members: 8 instructors verified dynamic availability slots.",
        "Administrators: 3 personnel tested account registration approvals."
    ]
    create_textbox_list(slide, Inches(7.1), Inches(2.7), Inches(5.4), Inches(3.7), meth2_bullets, font_size=13)

    # ----------------------------------------------------
    # SLIDE 16: SYSTEM EVALUATION RESULTS (Light Gray)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "SYSTEM EVALUATION RESULTS", "ISO/IEC 25010 Software Quality ratings")
    
    # Left Column: Metrics
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(4.5), Inches(2.3), C_WHITE, C_GRAY_LIGHT)
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(4.5), Inches(0.08), C_EMERALD)
    
    tb_res1 = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(4.1), Inches(1.9))
    tf_res1 = tb_res1.text_frame
    tf_res1.word_wrap = True
    p_res1 = tf_res1.paragraphs[0]
    p_res1.text = "Overall Evaluation Mean"
    p_res1.font.bold = True
    p_res1.font.size = Pt(13)
    p_res1.font.color.rgb = C_TEXT_MUTED
    p_res1.space_after = Pt(4)
    p_res1_2 = tf_res1.add_paragraph()
    p_res1_2.text = "4.71 / 5.00"
    p_res1_2.font.bold = True
    p_res1_2.font.size = Pt(36)
    p_res1_2.font.color.rgb = C_EMERALD
    p_res1_3 = tf_res1.add_paragraph()
    p_res1_3.text = "Interpreted as Excellent overall mean score, representing high quality rating."
    p_res1_3.font.size = Pt(11)
    p_res1_3.font.color.rgb = C_TEXT_DARK
    
    draw_rect(slide, Inches(0.5), Inches(4.4), Inches(4.5), Inches(2.4), C_WHITE, C_GRAY_LIGHT)
    draw_rect(slide, Inches(0.5), Inches(4.4), Inches(4.5), Inches(0.08), C_BLUE)
    
    tb_res2 = slide.shapes.add_textbox(Inches(0.7), Inches(4.6), Inches(4.1), Inches(2.0))
    tf_res2 = tb_res2.text_frame
    tf_res2.word_wrap = True
    p_res2 = tf_res2.paragraphs[0]
    p_res2.text = "Performance Efficiency"
    p_res2.font.bold = True
    p_res2.font.size = Pt(13)
    p_res2.font.color.rgb = C_TEXT_MUTED
    p_res2.space_after = Pt(4)
    p_res2_2 = tf_res2.add_paragraph()
    p_res2_2.text = "4.81 Rating"
    p_res2_2.font.bold = True
    p_res2_2.font.size = Pt(28)
    p_res2_2.font.color.rgb = C_BLUE
    p_res2_3 = tf_res2.add_paragraph()
    p_res2_3.text = "Highest score among criteria. Users highly satisfied with system speed, loading times, and real-time database updates."
    p_res2_3.font.size = Pt(11)
    p_res2_3.font.color.rgb = C_TEXT_DARK

    # Right Column: ISO/IEC 25010 Table
    tb_tbl_lbl = slide.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(7.3), Inches(0.4))
    tb_tbl_lbl.text_frame.paragraphs[0].text = "ISO/IEC 25010 Evaluation Results Breakdown"
    tb_tbl_lbl.text_frame.paragraphs[0].font.bold = True
    tb_tbl_lbl.text_frame.paragraphs[0].font.size = Pt(13.5)
    tb_tbl_lbl.text_frame.paragraphs[0].font.color.rgb = C_NAVY_LIGHT
    
    table_shape = slide.shapes.add_table(7, 3, Inches(5.5), Inches(2.3), Inches(7.3), Inches(4.5))
    tbl = table_shape.table
    tbl.columns[0].width = Inches(3.0)
    tbl.columns[1].width = Inches(2.0)
    tbl.columns[2].width = Inches(2.3)
    
    headers = ["Quality Characteristic", "Mean Rating", "Interpretation"]
    for col_idx, header in enumerate(headers):
        cell = tbl.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_NAVY_LIGHT
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        
    data = [
        ("Functionality", "4.72", "Excellent"),
        ("Usability", "4.65", "Excellent"),
        ("Reliability", "4.69", "Excellent"),
        ("Performance Efficiency", "4.81", "Excellent"),
        ("Security", "4.69", "Excellent"),
        ("Overall Mean", "4.71", "Excellent")
    ]
    
    for row_idx, row_data in enumerate(data):
        for col_idx, val in enumerate(row_data):
            cell = tbl.cell(row_idx + 1, col_idx)
            cell.text = val
            cell.fill.solid()
            if row_idx == 5:
                cell.fill.fore_color.rgb = RGBColor(237, 245, 255)
            else:
                cell.fill.fore_color.rgb = C_WHITE
            
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Arial"
            p.font.size = Pt(10.5)
            if row_idx == 5:
                p.font.bold = True
                p.font.color.rgb = C_NAVY_LIGHT
            else:
                p.font.color.rgb = C_TEXT_DARK

    # ----------------------------------------------------
    # SLIDE 17: DEVICE COMPATIBILITY (Light Gray)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "CROSS-DEVICE SCREEN COMPATIBILITY", "Viewport responsiveness tests across multiple devices")
    
    # Table layout
    table_shape = slide.shapes.add_table(6, 4, Inches(0.5), Inches(1.8), Inches(12.333), Inches(5.0))
    tbl = table_shape.table
    tbl.columns[0].width = Inches(2.2)
    tbl.columns[1].width = Inches(2.8)
    tbl.columns[2].width = Inches(2.3)
    tbl.columns[3].width = Inches(5.0)
    
    headers = ["Device Model", "Screen Size / Resolution", "Testing OS & Browser", "Display Suitability Status"]
    for col_idx, header in enumerate(headers):
        cell = tbl.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_NAVY_LIGHT
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        
    data = [
        ("Vivo Y17s", "6.56\" (720x1612 px)", "Android / Chrome Mobile", "PASSED (100% Mobile responsiveness)"),
        ("Samsung Galaxy S9", "5.8\" (1440x2960 px)", "Android / Chrome Mobile", "PASSED (100% Mobile touch scale)"),
        ("iPad Air Pro", "11\" (2360x1640 px)", "iOS / Chrome Tablet", "PASSED (100% Tablet layout double stack)"),
        ("MSI GF65 Thin", "15.6\" (1920x1080 px)", "Windows / Chrome Desktop", "PASSED (100% Laptop sidebar display)"),
        ("Desktop Computer", "24\" (1920x1080 px)", "Windows / Chrome Desktop", "PASSED (100% Desktop center bound)")
    ]
    
    for row_idx, row_data in enumerate(data):
        for col_idx, val in enumerate(row_data):
            cell = tbl.cell(row_idx + 1, col_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_WHITE
            
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Arial"
            p.font.size = Pt(10)
            p.font.color.rgb = C_TEXT_DARK
            if col_idx == 3:
                p.font.bold = True
                p.font.color.rgb = C_EMERALD

    # ----------------------------------------------------
    # SLIDE 18: CROSS-DEVICE RESPONSIVENESS MOCKUPS (Light Gray)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "CROSS-DEVICE RESPONSIVENESS MOCKUPS", "System layout adapting dynamically across different device viewports")
    
    # Left mockup: Mobile
    insert_screenshot(slide, "consultime_mobile_mockup.png", Inches(0.5), Inches(1.8), Inches(3.8), Inches(3.6))
    tb_rsp1 = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(3.8), Inches(1.7))
    tf_rsp1 = tb_rsp1.text_frame
    tf_rsp1.word_wrap = True
    p_rsp1 = tf_rsp1.paragraphs[0]
    p_rsp1.text = "Mobile Viewport. Optimizes touch layouts for small phone screens (e.g. Vivo Y17s & Samsung S9)."
    p_rsp1.font.name = "Arial"
    p_rsp1.font.size = Pt(11)
    p_rsp1.font.color.rgb = C_TEXT_MUTED
    p_rsp1.alignment = PP_ALIGN.CENTER
    
    # Center mockup: Tablet
    insert_screenshot(slide, "consultime_tablet_mockup.png", Inches(4.76), Inches(1.8), Inches(3.8), Inches(3.6))
    tb_rsp2 = slide.shapes.add_textbox(Inches(4.76), Inches(5.5), Inches(3.8), Inches(1.7))
    tf_rsp2 = tb_rsp2.text_frame
    tf_rsp2.word_wrap = True
    p_rsp2 = tf_rsp2.paragraphs[0]
    p_rsp2.text = "Tablet Viewport. Adapts sidebar controls and card grids for medium screens (e.g. iPad Air Pro)."
    p_rsp2.font.name = "Arial"
    p_rsp2.font.size = Pt(11)
    p_rsp2.font.color.rgb = C_TEXT_MUTED
    p_rsp2.alignment = PP_ALIGN.CENTER
    
    # Right mockup: Desktop
    insert_screenshot(slide, "consultime_desktop_mockup.png", Inches(9.02), Inches(1.8), Inches(3.8), Inches(3.6))
    tb_rsp3 = slide.shapes.add_textbox(Inches(9.02), Inches(5.5), Inches(3.8), Inches(1.7))
    tf_rsp3 = tb_rsp3.text_frame
    tf_rsp3.word_wrap = True
    p_rsp3 = tf_rsp3.paragraphs[0]
    p_rsp3.text = "Desktop Viewport. Maximizes display density and full telemetry views for laptop and desktop workstation resolutions."
    p_rsp3.font.name = "Arial"
    p_rsp3.font.size = Pt(11)
    p_rsp3.font.color.rgb = C_TEXT_MUTED
    p_rsp3.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # SLIDE 19: IMPLICATIONS & RECOMMENDATIONS (Light Gray)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_BG_LIGHT)
    add_slide_header(slide, "IMPLICATIONS FOR PRACTICE & RECOMMENDATIONS", "Avenues for academic digital transition")
    
    # Left Column: Implications
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), C_WHITE, C_GRAY_LIGHT)
    draw_rect(slide, Inches(0.5), Inches(1.8), Inches(6.0), Inches(0.12), C_NAVY_LIGHT)
    
    tb_imp = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(5.4), Inches(4.4))
    tf_imp = tb_imp.text_frame
    tf_imp.word_wrap = True
    p_ip = tf_imp.paragraphs[0]
    p_ip.text = "Implications for Practice"
    p_ip.font.bold = True
    p_ip.font.size = Pt(16)
    p_ip.font.color.rgb = C_NAVY_LIGHT
    p_ip.space_after = Pt(14)
    
    imp_bullets = [
        "Academic Efficiency: Shifting physical appointment grids online removes scheduling conflicts.",
        "Role Isolation Security: Enforcing strict Row Level Security rules sets a baseline design standard for storing student and advisor logs securely."
    ]
    create_textbox_list(slide, Inches(0.8), Inches(2.7), Inches(5.4), Inches(3.7), imp_bullets, font_size=12.5)

    # Right Column: Recommendations
    draw_rect(slide, Inches(6.8), Inches(1.8), Inches(6.0), Inches(5.0), C_WHITE, C_GRAY_LIGHT)
    draw_rect(slide, Inches(6.8), Inches(1.8), Inches(6.0), Inches(0.12), C_EMERALD)
    
    tb_rec = slide.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.4), Inches(4.4))
    tf_rec = tb_rec.text_frame
    tf_rec.word_wrap = True
    p_rc = tf_rec.paragraphs[0]
    p_rc.text = "Recommendations for Future Studies"
    p_rc.font.bold = True
    p_rc.font.size = Pt(16)
    p_rc.font.color.rgb = C_NAVY_LIGHT
    p_rc.space_after = Pt(14)
    
    rec_bullets = [
        "Live Chat Channels: Integrate context-based chatting views adjacent to booking slots.",
        "Cross-Platform Deployment: Expand Capacitor pipelines to package applications for Apple iOS systems.",
        "Video Conferences: Embed direct browser video API bridges for remote consultation sessions."
    ]
    create_textbox_list(slide, Inches(7.1), Inches(2.7), Inches(5.4), Inches(3.7), rec_bullets, font_size=12.5)

    # ----------------------------------------------------
    # SLIDE 20: THANK YOU (Dark Navy)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, C_NAVY_DARK)
    
    # Background glowing accents
    draw_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), C_BLUE)
    draw_rect(slide, Inches(13.18), Inches(0), Inches(0.15), Inches(7.5), C_EMERALD)
    
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "THANK YOU!"
    p1.font.name = "Arial"
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = C_MINT
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(14)
    
    p2 = tf.add_paragraph()
    p2.text = "ConsulTime: A Cross-Platform Consultation and Appointment Management System"
    p2.font.name = "Arial"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_after = Pt(36)
    
    p3 = tf.add_paragraph()
    p3.text = "Questions & Recommendations Session"
    p3.font.name = "Arial"
    p3.font.size = Pt(16)
    p3.font.bold = True
    p3.font.color.rgb = C_BLUE
    p3.alignment = PP_ALIGN.CENTER
    
    p4 = tf.add_paragraph()
    p4.text = "Philipp Edward Sapalicio, Ednell Sapalicio"
    p4.font.name = "Arial"
    p4.font.size = Pt(13)
    p4.font.color.rgb = C_TEXT_MUTED
    p4.alignment = PP_ALIGN.CENTER
    p4.space_before = Pt(20)

    # Save to workspace target filename
    try:
        prs.save("ConsulTime Mahayag Presentation.pptx")
        print("SUCCESS: Programmatically compiled ConsulTime Mahayag Presentation.pptx!")
    except PermissionError:
        alternative_name = "ConsulTime Mahayag Presentation (Updated).pptx"
        prs.save(alternative_name)
        print(f"WARNING: 'ConsulTime Mahayag Presentation.pptx' is locked (probably open in PowerPoint).")
        print(f"SUCCESS: Saved compilation to alternative path: {alternative_name}")

if __name__ == "__main__":
    create_presentation()
